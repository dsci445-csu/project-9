from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

p = Presentation()

title_slide = p.slides.add_slide(p.slide_layouts[0])
title_slide.placeholders[0].text = "LoL Match Review Tool"
title_slide.placeholders[1].text = "Ryan Stevens, Nick Lee, Elle Angell, Cole Garnier"






p.save("presentation.pptx")

def get_layouts():
    p = Presentation()

    for i in range (len(p.slide_layouts)):
        s = p.slides.add_slide(p.slide_layouts[i])
        if s.shapes.title is not None:
            s.shapes.title.text = f"This is slide no. {i}"
        if s.placeholders is not None:
            for j in range(len(s.placeholders)):
                s.placeholders[j].text = f"{i} placeholder {j}"

    p.save('layouts.pptx')
